import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Use widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    BG_COLOR = RGBColor(15, 16, 19)        # #0F1013
    CARD_BG = RGBColor(31, 32, 38)         # #1F2026
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_LIGHT_GRAY = RGBColor(200, 200, 205)
    TEXT_MUTED = RGBColor(150, 150, 155)
    
    ACCENT_GREEN = RGBColor(57, 255, 20)   # #39FF14
    ACCENT_BLUE = RGBColor(0, 240, 255)    # #00F0FF
    ACCENT_RED = RGBColor(255, 59, 48)     # #FF3B30
    BORDER_COLOR = RGBColor(50, 52, 62)
    
    blank_layout = prs.slide_layouts[6]
    
    def set_slide_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_title(slide, text, accent_color=ACCENT_BLUE):
        # Vertical accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.55), Inches(0.08), Inches(0.5)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_color
        bar.line.fill.background()
        
        # Title textbox
        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(0.45), Inches(11.5), Inches(0.7))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

    def add_subtitle(slide, text):
        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.1), Inches(11.5), Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_LIGHT_GRAY

    def draw_card(slide, left, top, width, height, fill_color=CARD_BG, border_color=BORDER_COLOR):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()
        return shape

    def add_card_text(slide, left, top, width, height, title, body_lines, accent_color=None, title_size=18, body_size=12):
        # Textbox inside card
        txBox = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.25), width - Inches(0.5), height - Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_left = Inches(0)
        
        # Card Title
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.name = 'Microsoft YaHei'
        p_title.font.size = Pt(title_size)
        p_title.font.bold = True
        p_title.font.color.rgb = accent_color if accent_color else TEXT_WHITE
        p_title.space_after = Pt(12)
        
        # Card Body
        for line in body_lines:
            p = tf.add_paragraph()
            p.text = line
            p.font.name = 'Microsoft YaHei'
            p.font.size = Pt(body_size)
            p.font.color.rgb = TEXT_LIGHT_GRAY
            p.space_after = Pt(6)
            
            # If line starts with specific bullet chars, we can style it
            if line.strip().startswith("•") or line.strip().startswith("-"):
                p.level = 0
                
    # ==================== SLIDE 1: COVER ====================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide1)
    
    # Decorative shape - large subtle colored background glow
    glow = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(1), Inches(4), Inches(4))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(20, 40, 50)
    glow.line.fill.background()
    
    # Title Text Frame
    txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    
    p0 = tf.paragraphs[0]
    p0.text = "Flying Run"
    p0.font.name = 'Arial'
    p0.font.size = Pt(64)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_BLUE
    p0.space_after = Pt(10)
    
    p1 = tf.add_paragraph()
    p1.text = "端云协同的AI陪伴式运动打卡与健康数据监测系统"
    p1.font.name = 'Microsoft YaHei'
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.space_after = Pt(15)
    
    p2 = tf.add_paragraph()
    p2.text = "面向校园隐私安全、社区适老化康养与个人高情感伴跑的多模态智能健康平台"
    p2.font.name = 'Microsoft YaHei'
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_LIGHT_GRAY
    p2.space_after = Pt(40)
    
    p3 = tf.add_paragraph()
    p3.text = "【 创新工程实践报告 & 商业计划书 】"
    p3.font.name = 'Microsoft YaHei'
    p3.font.size = Pt(18)
    p3.font.bold = True
    p3.font.color.rgb = ACCENT_GREEN

    # ==================== SLIDE 2: AGENDA ====================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide2)
    add_title(slide2, "目录 / 汇报大纲", ACCENT_BLUE)
    add_subtitle(slide2, "Flying Run 全方位汇报内容框架")
    
    agenda_items = [
        ("01", "项目背景与市场契机", "政策导向与数字体育市场爆发"),
        ("02", "行业三大核心痛点", "C端/B端/G端痛点与真实调研"),
        ("03", "产品核心解决方案", "三大矩阵商业与技术模式"),
        ("04", "端云协同技术壁垒", "算法微调、情感计算与架构设计"),
        ("05", "用户体验与游戏化", "玻璃美学设计与八角游戏化框架"),
        ("06", "商业模式与财务预测", "盈利矩阵、极低边际成本与里程碑")
    ]
    
    for idx, (num, title, desc) in enumerate(agenda_items):
        col = idx % 3
        row = idx // 3
        
        left = Inches(1.0 + col * 3.8)
        top = Inches(2.2 + row * 2.2)
        width = Inches(3.5)
        height = Inches(1.8)
        
        draw_card(slide2, left, top, width, height)
        
        # Textbox inside card
        txBox = slide2.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_left = Inches(0)
        
        # Number
        p_num = tf.paragraphs[0]
        p_num.text = num
        p_num.font.name = 'Arial'
        p_num.font.size = Pt(22)
        p_num.font.bold = True
        p_num.font.color.rgb = ACCENT_GREEN
        p_num.space_after = Pt(2)
        
        # Title
        p_title = tf.add_paragraph()
        p_title.text = title
        p_title.font.name = 'Microsoft YaHei'
        p_title.font.size = Pt(15)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE
        p_title.space_after = Pt(4)
        
        # Desc
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = 'Microsoft YaHei'
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = TEXT_MUTED

    # ==================== SLIDE 3: BACKGROUND ====================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide3)
    add_title(slide3, "一、 政策依据与数字体育市场爆发", ACCENT_BLUE)
    add_subtitle(slide3, "全民健康战略与跨部门协作机制，驱动运动促进健康事业的高质量发展")
    
    # 2 columns layout
    # Left Column: Policy
    left_x = Inches(0.8)
    top_y = Inches(1.8)
    w_card = Inches(5.6)
    h_card = Inches(4.8)
    
    draw_card(slide3, left_x, top_y, w_card, h_card)
    add_card_text(
        slide3, left_x, top_y, w_card, h_card,
        "政策支撑：国家战略关口前移",
        [
            "• 《健康中国2030规划纲要》",
            "  - 明确将全民健身作为国家战略，目标到2030年经常参加体育锻炼的人数比例达40%以上。",
            "",
            "• 《关于推动运动促进健康事业高质量发展的指导意见》(2025)",
            "  - 明确“体医融合”导向，提倡在社区推广科学的运动干预，打通体育、卫健、民政及教育部门协作。",
            "",
            "• 十五五群体工作规划",
            "  - 推动数字化转型，以“数字体育大脑”、“运动码”试点打破公共服务不均等，实现治理现代化。"
        ],
        ACCENT_BLUE, title_size=18, body_size=13
    )
    
    # Right Column: Market Growth
    right_x = Inches(6.8)
    draw_card(slide3, right_x, top_y, w_card, h_card)
    add_card_text(
        slide3, right_x, top_y, w_card, h_card,
        "市场红利：大健康与AI陪伴双重驱动",
        [
            "• AI 智能健身软件市场",
            "  - 2024年国内规模达523亿元，预计2025年激增至821.6亿元，活跃用户接近3亿。",
            "",
            "• 智能穿戴设备市场",
            "  - IDC预测2025年国内智能穿戴大盘将达707亿元，年复合增长率（CAGR）维持在14.9%以上。",
            "",
            "• AI 情感陪伴赛道",
            "  - 2025年全球规模约377亿美元，预计2034年将激增至4300亿美元，年复合增长率达31.3%。",
            "  - 本系统首创AI陪伴与运动健康融合，深度切入情感刚需。"
        ],
        ACCENT_GREEN, title_size=18, body_size=13
    )

    # ==================== SLIDE 4: PAIN POINTS ====================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide4)
    add_title(slide4, "二、 痛点分析：数字体育与监测生态错位", ACCENT_RED)
    add_subtitle(slide4, "当前C端、B端校园、G端社区健康产品均存在严重未满足的强烈需求")
    
    # 3 columns layout
    col_w = Inches(3.6)
    col_h = Inches(4.8)
    col_top = Inches(1.8)
    
    pains = [
        ("C端运动应用 (如 Keep、悦跑圈)", 
         ["• 商业化吃相难看", "  - 原本单纯的记录工具沦为“广告+拼多多+小红书”，基础数据导出也被VIP门槛卡死。", "", "• 孤独感与情绪匮乏", "  - 死板的机器音缺乏情感，凡尔赛式的秀场社区造成普通用户流失。30日留存率出现断崖式下跌。"], ACCENT_RED),
        ("校园阳光跑 (如 步道乐跑等)", 
         ["• “技术霸权”引发逆反", "  - GPS定位频繁漂移、后台被杀丢数据，导致学生对强制跑步产生极强抵触情绪。", "", "• 黑产泛滥与隐私红线", "  - 催生了代跑、定位作弊灰产；强制收集地理轨迹和身份隐私，极易面临数据合规风险。"], ACCENT_RED),
        ("社区适老化监测 (适老硬件/App)", 
         ["• “伪适老”与高误报率", "  - 大多仅放大字体，层级交互依旧复杂；防跌倒监测误报率极高（如老人用力拍腿即报警）。", "", "• “被监视”与充电焦虑", "  - 监控设备剥夺老人尊严，一天一充严重损害便利性，老人抗拒使用导致功能变伪需求。"], ACCENT_RED)
    ]
    
    for i, (title, lines, color) in enumerate(pains):
        left_pos = Inches(0.8 + i * 4.0)
        draw_card(slide4, left_pos, col_top, col_w, col_h)
        add_card_text(slide4, left_pos, col_top, col_w, col_h, title, lines, color, title_size=16, body_size=12)

    # ==================== SLIDE 5: SOLUTIONS ====================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide5)
    add_title(slide5, "三、 核心解决方案：三大矩阵模式", ACCENT_GREEN)
    add_subtitle(slide5, "通过高情感温度、高隐私安全、高算法精度的服务体系，精准覆盖细分人群")
    
    solutions = [
        ("AI陪伴运动伴跑 (B2C大众跑者)", 
         ["• GPT-SoVITS 人声复刻", "  - 上传5秒音频即可生成恋人、亲人或动漫角色音色，实时在跑步中伴跑。", "", "• 生理-情感共鸣计算", "  - 根据心率自动切换语调：平稳时闲聊，心率过高时警示，终点前高亢激昂。", "", "• 千人千面运动处方", "  - 汇总体征数据，生成量身定制的运动和饮食规划建议。"], ACCENT_GREEN),
        ("校园隐私打卡 (B2B2C学校教务)", 
         ["• “数据不落地”代理模式", "  - App作为富客户端代理，将轨迹等敏感信息直接加密直传校方内网，开发商绝不留存。", "", "• 多维传感器交叉验证", "  - 支持穿戴手环与GPS比对，校验心率/步频与定位轨迹的一致性，智能判定代跑等作弊行为。"], ACCENT_GREEN),
        ("社区康养静默守护 (B2G/B2B康养)", 
         ["• 后台无感静默守护", "  - 运行时无感驻留后台，不需要老年人操作App或进行任何多余交互。", "", "• 云端大模型智能降噪", "  - 结合原始传感器信号与心率突变二次校验，将挥手等跌倒误报率大幅降低85%以上。", "", "• 一键连线全科医生", "  - 异常时联动Direct Dial，自动报警、发送定位与推送。"], ACCENT_GREEN)
    ]
    
    for i, (title, lines, color) in enumerate(solutions):
        left_pos = Inches(0.8 + i * 4.0)
        draw_card(slide5, left_pos, col_top, col_w, col_h)
        add_card_text(slide5, left_pos, col_top, col_w, col_h, title, lines, color, title_size=16, body_size=12)

    # ==================== SLIDE 6: TECHNICAL MOAT ====================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide6)
    add_title(slide6, "四、 端云协同架构与极低边际成本", ACCENT_BLUE)
    add_subtitle(slide6, "摆脱对重度闭源大模型的依赖，通过开源模型本地化微调构筑技术与成本护城河")
    
    draw_card(slide6, left_x, top_y, w_card, h_card)
    add_card_text(
        slide6, left_x, top_y, w_card, h_card,
        "技术架构壁垒",
        [
            "• 跨平台富前端架构",
            "  - 采用 Flutter + Dart 构建，通过 Adapter 适配器模式，搭积木式零成本接入华为、小米等主流穿戴生态。",
            "",
            "• 云端极速声音克隆管道",
            "  - 基于 GPT-SoVITS 在云端 GPU 实时加速部署，实现低延迟语音合成与动态伴跑流推送。",
            "",
            "• 数据解耦与加密网关",
            "  - 采用 Dio 拦截网关结合本地国密算法，在前端直接完成学籍信息与运动轨迹的组装和加密传输。"
        ],
        ACCENT_BLUE, title_size=18, body_size=13
    )
    
    draw_card(slide6, right_x, top_y, w_card, h_card)
    add_card_text(
        slide6, right_x, top_y, w_card, h_card,
        "极低边际成本结构",
        [
            "• 大模型微调极低开销 (<5万人民币)",
            "  - 弃用昂贵的闭源GPT大模型，本地微调开源 Qwen3-7B。引入 QLoRA 与 Paged Optimizers，单张消费级显卡(如4090)即可完成医学偏好对齐训练。",
            "",
            "• 声音克隆低运行成本 (<1万人民币)",
            "  - 少样本克隆仅需 5 秒录音。推理采用轻量级异步请求，随着用户规模上升，服务器带宽及计算开销平缓。",
            "",
            "• 边际推理成本接近于零",
            "  - 避免API按Token计费模式，单次训练完成后，私有云部署可承受极高并发量。"
        ],
        ACCENT_GREEN, title_size=18, body_size=13
    )

    # ==================== SLIDE 7: ALGORITHM DPO ====================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide7)
    add_title(slide7, "五、 算法创新：QLoRA + Qwen3-7B + DPO 运动处方", ACCENT_BLUE)
    add_subtitle(slide7, "依据2025年《临床运动处方实践专家共识》进行DPO医学专家偏好对齐，解决死板硬编码处方问题")
    
    # Left description card
    left_w = Inches(4.5)
    draw_card(slide7, Inches(0.8), Inches(1.8), left_w, Inches(4.8))
    add_card_text(
        slide7, Inches(0.8), Inches(1.8), left_w, Inches(4.8),
        "偏好对齐原理",
        [
            "• 彻底抛弃硬编码静态逻辑",
            "  - 传统APP只能依靠预设公式输出简陋模版，本系统大模型能综合评估多维体征。",
            "",
            "• 循证医学标准融合",
            "  - 针对运动风险筛查、安全心率区间、特殊病理进行医学奖励打分。",
            "",
            "• DPO 偏好对齐",
            "  - 通过 Direct Preference Optimization，训练模型规避“过度激进、不顾疲劳”的推荐，转向“符合医学安全共识”的主动降级处方。"
        ],
        ACCENT_BLUE, title_size=18, body_size=12
    )
    
    # Right code/table card
    right_w = Inches(6.8)
    draw_card(slide7, Inches(5.7), Inches(1.8), right_w, Inches(4.8))
    
    # Custom table inside the right card for Chosen/Rejected comparison
    txBox = slide7.shapes.add_textbox(Inches(5.95), Inches(2.05), right_w - Inches(0.5), Inches(4.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0)
    tf.margin_left = Inches(0)
    
    p_t = tf.paragraphs[0]
    p_t.text = "医学 DPO 对齐典型样本设计对比"
    p_t.font.name = 'Microsoft YaHei'
    p_t.font.size = Pt(16)
    p_t.font.bold = True
    p_t.font.color.rgb = ACCENT_GREEN
    p_t.space_after = Pt(14)
    
    p_in = tf.add_paragraph()
    p_in.text = "输入：48岁男性，昨晚睡眠4.2h（极低），今日静息心率 84bpm（偏高），诉求：高强度跑步"
    p_in.font.name = 'Microsoft YaHei'
    p_in.font.size = Pt(11)
    p_in.font.bold = True
    p_in.font.color.rgb = TEXT_WHITE
    p_in.space_after = Pt(8)
    
    # Chosen block
    p_c_title = tf.add_paragraph()
    p_c_title.text = "✓ 推荐(Chosen) - 符合《共识》医学干预要求"
    p_c_title.font.name = 'Microsoft YaHei'
    p_c_title.font.size = Pt(12)
    p_c_title.font.bold = True
    p_c_title.font.color.rgb = ACCENT_GREEN
    
    p_c_val = tf.add_paragraph()
    p_c_val.text = "“因您昨晚睡眠过少且心率偏高，有心肌超负荷风险。今日不建议原定燃脂跑，今日处方调整为：30分钟低强度慢走（LSD），靶心率控制在110-125 bpm，侧重身体恢复。”"
    p_c_val.font.name = 'Microsoft YaHei'
    p_c_val.font.size = Pt(11)
    p_c_val.font.color.rgb = TEXT_LIGHT_GRAY
    p_c_val.space_after = Pt(12)
    
    # Rejected block
    p_r_title = tf.add_paragraph()
    p_r_title.text = "✗ 拒绝(Rejected) - 死板硬套或违背医学常识"
    p_r_title.font.name = 'Microsoft YaHei'
    p_r_title.font.size = Pt(12)
    p_r_title.font.bold = True
    p_r_title.font.color.rgb = ACCENT_RED
    
    p_r_val = tf.add_paragraph()
    p_r_val.text = "“好的，已为您开启燃脂跑模式。今天请坚持跑完 5 公里，配速 5 分钟，加油坚持就是胜利！”(容易诱发心梗等突发风险)"
    p_r_val.font.name = 'Microsoft YaHei'
    p_r_val.font.size = Pt(11)
    p_r_val.font.color.rgb = TEXT_MUTED

    # ==================== SLIDE 8: AFFECTIVE COMPUTING ====================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide8)
    add_title(slide8, "六、 陪伴创新：GPT-SoVITS声音克隆与情感计算", ACCENT_GREEN)
    add_subtitle(slide8, "5秒声音克隆结合生理体征解耦情感计算，让AI伴跑教练真正拥有生命质感")
    
    # 2 columns layout
    draw_card(slide8, left_x, top_y, w_card, h_card)
    add_card_text(
        slide8, left_x, top_y, w_card, h_card,
        "5秒极速声音复刻机制",
        [
            "• 极低门槛的音色克隆",
            "  - 用户仅需录制或上传5秒的亲人、恋人或偶像音频片段，云端 GPT-SoVITS 就能高保真复刻其音色、情绪与语调起伏。",
            "",
            "• 跨时间的情感记忆保留",
            "  - AI 教练会记住用户上一周的健康抱怨(如“膝盖微酸”)，在本次运动开始前主动进行关切询问，建立长期信任关系纽带。",
            "",
            "• 规避外界纪律约束焦虑",
            "  - 将枯燥的运动打卡，转化为与“亲密关系”共同度过的温暖时光，极大提升用户粘性与打卡动力。"
        ],
        ACCENT_BLUE, title_size=18, body_size=13
    )
    
    draw_card(slide8, right_x, top_y, w_card, h_card)
    add_card_text(
        slide8, right_x, top_y, w_card, h_card,
        "运动生理信号的情感解耦",
        [
            "• 多维生理特征情感计算",
            "  - 运动中的心率变异性(HRV)、步频呼吸比及加速度波形，是身体极限与情绪精力的真实映射。",
            "",
            "• 伴跑语速与语调动态生成",
            "  - 心率平稳时(<120 bpm): 伴跑人声温和，以闲聊语气科普健康知识，给予陪伴。",
            "  - 达到上限时(>165 bpm): 语气紧急急促，发出主动警示(如“心率太高了，听我的，快慢下来！”)。",
            "  - 冲刺阶段: 语气高亢激昂，释放强力激励。"
        ],
        ACCENT_GREEN, title_size=18, body_size=13
    )

    # ==================== SLIDE 9: PRIVACY GATEWAY ====================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide9)
    add_title(slide9, "七、 架构创新：“数据不落地”的隐私打卡", ACCENT_BLUE)
    add_subtitle(slide9, "首创富前端交互代理架构，将运动打卡与学校/社区医院内网直连，打破商业隐私顾虑")
    
    # Draw simple flow elements inside cards
    draw_card(slide9, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8))
    add_card_text(
        slide9, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8),
        "1. 终端感知与清洗",
        [
            "• 归一化手环接口",
            "  - 手机APP通过蓝牙拉取华为、小米手环或手机内置传感器的步数、心率与GPS轨迹。",
            "",
            "• 脱敏上报云端",
            "  - 过滤掉个人姓名、学号等敏感身份信息，仅将纯粹生理指标上传至 Flying Run 算法云进行处方与音频推理。"
        ],
        ACCENT_BLUE, title_size=16, body_size=12
    )
    
    draw_card(slide9, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8))
    add_card_text(
        slide9, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8),
        "2. 数据不落地中转",
        [
            "• 纯代理中转网关",
            "  - Flying Run 服务端作为中转，不建立任何轨迹数据库，不滞留个人隐私行踪。",
            "",
            "• 加密直传校园内网",
            "  - 轨迹与学籍信息在本地组装，以高强度非对称加密直接回传校方教务系统，完全规避泄露红线。"
        ],
        ACCENT_GREEN, title_size=16, body_size=12
    )
    
    draw_card(slide9, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.8))
    add_card_text(
        slide9, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.8),
        "3. 安全防作弊校验",
        [
            "• 多维传感器交叉验证",
            "  - 传统代跑/模拟器定位作弊仅篡改GPS轨迹，本系统提取智能硬件获取的心率/步频波动曲线。",
            "",
            "• 动态逻辑鉴权",
            "  - 智能算法检测到轨迹变动但手环心率与步数未发生响应波动时，自动在后台判定为疑似作弊，确保考核公平性。"
        ],
        ACCENT_BLUE, title_size=16, body_size=12
    )

    # ==================== SLIDE 10: UI/UX EXPERIENCE ====================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide10)
    add_title(slide10, "八、 用户体验：暗色玻璃微光美学与灵动交互", ACCENT_BLUE)
    add_subtitle(slide10, "高端科技质感UI，配合流畅阻尼动画，构建高情感化、无干扰的运动交互界面")
    
    draw_card(slide10, left_x, top_y, w_card, h_card)
    add_card_text(
        slide10, left_x, top_y, w_card, h_card,
        "暗色玻璃拟态 (Glassmorphism)",
        [
            "• 极客暗黑主题底色",
            "  - 采用 `#0F1013` 深色背景，大幅降低在运动、强光或夜跑时的视觉刺眼感，提供高度沉浸感。",
            "",
            "• 炫彩霓虹配色与拟物卡片",
            "  - 卡片使用 `#1F2026` 配合背景高斯模糊(`blur(20px)`)及发光边框，搭配动力绿、科技蓝与警示红点缀。",
            "",
            "• 科技感字体系规范",
            "  - 引入 Outfit 现代字体系，以及核心数据展示上的 Orbitron 科技数码字体，极具运动张力与科技美感。"
        ],
        ACCENT_BLUE, title_size=18, body_size=13
    )
    
    draw_card(slide10, right_x, top_y, w_card, h_card)
    add_card_text(
        slide10, right_x, top_y, w_card, h_card,
        "灵动岛组件与动态指示器",
        [
            "• 全局轻量级灵动岛(Dynamic Island)",
            "  - 运动时有伴跑语音、打卡点触发或生理预警时，顶端灵动岛以平滑阻尼动画展开，不打断主流程操作。",
            "",
            "• 动态心率呼吸灯 (PulsingDot)",
            "  - 心率卡片指示灯会根据手环回传的实时心率调整其闪烁频率，赋予应用生命跳动质感。",
            "",
            "• 拟音轨波形组件",
            "  - AI 伴跑教练在发声时，界面下方会渲染实时根据音频频谱跳动的迷你音轨柱，增强人机交互反馈。"
        ],
        ACCENT_GREEN, title_size=18, body_size=13
    )

    # ==================== SLIDE 11: GAMIFICATION ====================
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide11)
    add_title(slide11, "九、 游戏化驱动与弹性健康打卡", ACCENT_GREEN)
    add_subtitle(slide11, "应用八角游戏化框架，变“外界强制纪律约束”为“内在心理游戏化驱动”")
    
    draw_card(slide11, left_x, top_y, w_card, h_card)
    add_card_text(
        slide11, left_x, top_y, w_card, h_card,
        "八角游戏化机制应用",
        [
            "• 史诗意义与使命感 (Epic Meaning)",
            "  - 跑步里程转换为“绿色碳减排点数”，汇入高校/班级榜，达标时由企业出资以班级名义在西北种下真树。",
            "",
            "• 解锁发音人声勋章 (Development)",
            "  - 连续打卡解锁特定 AI 音色参数(例如元气傲娇音、魔鬼严格教练音)及炫酷三维勋章卡片。",
            "",
            "• 社交关联与宿舍小镇 (Social Influence)",
            "  - 情侣可跨地域声音克隆伴跑；宿舍组队建设虚拟小镇，任何一人断跑会导致小镇停电，对抗个人惰性。"
        ],
        ACCENT_BLUE, title_size=18, body_size=13
    )
    
    draw_card(slide11, right_x, top_y, w_card, h_card)
    add_card_text(
        slide11, right_x, top_y, w_card, h_card,
        "弹性健康任务分配逻辑",
        [
            "• 智能拉取手环最新指标",
            "  - 当用户触发跑步打卡，系统立即检查睡眠质量与静息心率是否处于异常区间。",
            "",
            "• 生理保护主动任务降级",
            "  - 判定睡眠<5h或静息心率>85bpm时，大模型评估其处于疲劳状态。自动将“3公里跑”降级为“1.5公里漫步”。",
            "",
            "• 天气自适应推荐与断连保护",
            "  - 检测到雾霾或暴雨时，自动切换为室内徒步/跳绳打卡。降级打卡仍算作连续，保护用户成就感。"
        ],
        ACCENT_GREEN, title_size=18, body_size=13
    )

    # ==================== SLIDE 12: BUSINESS MODEL ====================
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide12)
    add_title(slide12, "十、 商业模式与三端盈利矩阵", ACCENT_BLUE)
    add_subtitle(slide12, "打通学校B端、社区G端与大众跑者C端，建立健康促进生态闭环")
    
    # 3 columns layout
    draw_card(slide12, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8))
    add_card_text(
        slide12, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8),
        "B端：智慧校园服务",
        [
            "• 目标群体",
            "  - 大中专院校体育教务部门。",
            "",
            "• 盈利模式",
            "  - 提供防作弊校园跑系统 SaaS 平台订阅费；",
            "  - 按年/按学生人数收取数据网关中转与本地系统维保费。",
            "",
            "• 客单价估算",
            "  - 平均每所高校收取约 **5万元 / 年** 的服务费。"
        ],
        ACCENT_BLUE, title_size=16, body_size=12
    )
    
    draw_card(slide12, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8))
    add_card_text(
        slide12, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8),
        "G端：智慧社区康养",
        [
            "• 目标群体",
            "  - 街道办、社区卫生中心、适老康养机构。",
            "",
            "• 盈利模式",
            "  - 老人健康监护 API 接口开发费；",
            "  - 社区物联网硬件中继网关部署及维护；",
            "  - 按年收取远程高危预警与医生连线维保费。",
            "",
            "• 客单价估算",
            "  - 平均每个社区点位收取约 **3万元 / 年**。"
        ],
        ACCENT_GREEN, title_size=16, body_size=12
    )
    
    draw_card(slide12, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.8))
    add_card_text(
        slide12, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.8),
        "C端：增值与分销变现",
        [
            "• 声音克隆增值服务",
            "  - 基础克隆免费1个，增加音色或扩充云存储空间收取存储费 (**15元/月**)。",
            "",
            "• 大模型医学Token包",
            "  - 超额医学处方出具，按使用Token计费 (**20元/10万字**)。",
            "",
            "• 硬件分销 CPS 佣金",
            "  - 推荐未购硬件用户购买合作手环，抽取 **8%-12%** 佣金。"
        ],
        ACCENT_BLUE, title_size=16, body_size=12
    )

    # ==================== SLIDE 13: FINANCIALS ====================
    slide13 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide13)
    add_title(slide13, "十一、 财务预测与融资里程碑规划", ACCENT_BLUE)
    add_subtitle(slide13, "三年期用户规模与盈利预测，及天使轮融资资金用途规划")
    
    # Left column: Table of 3 years
    draw_card(slide13, left_x, top_y, w_card, h_card)
    
    txBox = slide13.shapes.add_textbox(left_x + Inches(0.25), top_y + Inches(0.25), w_card - Inches(0.5), h_card - Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0)
    tf.margin_left = Inches(0)
    
    p_tab = tf.paragraphs[0]
    p_tab.text = "三年期运营与财务推演指标"
    p_tab.font.name = 'Microsoft YaHei'
    p_tab.font.size = Pt(18)
    p_tab.font.bold = True
    p_tab.font.color.rgb = ACCENT_BLUE
    p_tab.space_after = Pt(12)
    
    metrics = [
        ("第一年 (Month 1-12):", "  - C端用户达 10万人，签约高校3所，社区10个。", "  - 预计年度总营业收入：**85 万元**，净利润率 -15%。"),
        ("第二年 (Month 13-24):", "  - C端用户达 80万人，签约高校35所，社区60个。", "  - 预计年度总营业收入：**820 万元**，净利润率 +12%。"),
        ("第三年 (Month 25-36):", "  - C端用户达 500万人，签约高校120所，社区350个。", "  - 预计年度总营业收入：**2640 万元**，净利润率 +36%。")
    ]
    
    for yr_t, user_m, fin_m in metrics:
        p_y = tf.add_paragraph()
        p_y.text = yr_t
        p_y.font.name = 'Microsoft YaHei'
        p_y.font.size = Pt(12)
        p_y.font.bold = True
        p_y.font.color.rgb = TEXT_WHITE
        p_y.space_after = Pt(2)
        
        p_u = tf.add_paragraph()
        p_u.text = user_m
        p_u.font.name = 'Microsoft YaHei'
        p_u.font.size = Pt(11)
        p_u.font.color.rgb = TEXT_LIGHT_GRAY
        p_u.space_after = Pt(2)
        
        p_f = tf.add_paragraph()
        p_f.text = fin_m
        p_f.font.name = 'Microsoft YaHei'
        p_f.font.size = Pt(11)
        p_f.font.color.rgb = ACCENT_GREEN
        p_f.space_after = Pt(8)
        
    # Right column: Funding & Milestones
    draw_card(slide13, right_x, top_y, w_card, h_card)
    add_card_text(
        slide13, right_x, top_y, w_card, h_card,
        "天使轮融资规划：150万 - 200万元",
        [
            "• 云算力与医学微调开支 (占比 40% / 约80万元)",
            "  - 用于购买 DPO 临床专家微调数据集，云服务器 GPU 算力租赁，以及 GPT-SoVITS 部署费。",
            "",
            "• 前后端研发与团队开支 (占比 35% / 约70万元)",
            "  - 扩充 Flutter 开发团队，完成多厂商智能穿戴设备 Health Kit SDK 的适配器(Adapter)组件开发。",
            "",
            "• 高校与社区卫生中心地推 (占比 25% / 约50万元)",
            "  - 设立首批校园及社区示范标杆点，打通教务内网路由，获取品牌背书与冷启动用户量。"
        ],
        ACCENT_GREEN, title_size=18, body_size=12
    )

    # ==================== SLIDE 14: THANK YOU ====================
    slide14 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide14)
    
    # Large colored background glow
    glow2 = slide14.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2), Inches(1.5), Inches(9), Inches(4.5))
    glow2.fill.solid()
    glow2.fill.fore_color.rgb = RGBColor(18, 30, 42)
    glow2.line.fill.background()
    
    txBox = slide14.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    
    p0 = tf.paragraphs[0]
    p0.text = "Flying Run"
    p0.font.name = 'Arial'
    p0.font.size = Pt(60)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_BLUE
    p0.alignment = PP_ALIGN.CENTER
    p0.space_after = Pt(12)
    
    p1 = tf.add_paragraph()
    p1.text = "谢谢观看，敬请各位评委与投资人指正！"
    p1.font.name = 'Microsoft YaHei'
    p1.font.size = Pt(24)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.alignment = PP_ALIGN.CENTER
    p1.space_after = Pt(15)
    
    p2 = tf.add_paragraph()
    p2.text = "端云协同的AI陪伴式运动打卡与健康数据监测系统"
    p2.font.name = 'Microsoft YaHei'
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_LIGHT_GRAY
    p2.alignment = PP_ALIGN.CENTER
    
    prs.save("Flying_Run_Business_Presentation.pptx")
    print("PPTX saved successfully to e:\\Flying_run\\Flying_Run_Business_Presentation.pptx")

if __name__ == '__main__':
    create_presentation()
